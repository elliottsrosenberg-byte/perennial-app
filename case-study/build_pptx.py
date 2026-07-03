#!/usr/bin/env python3
"""Editable Perennial walkthrough deck (.pptx) — document/resume style, GT America,
black-on-white, one centered mockup per slide, full-sentence captions."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(BASE, "screenshots")
OUT = os.path.join(BASE, "Perennial-Walkthrough.pptx")

INK   = RGBColor(0x1a, 0x1a, 0x1a)
SOFT  = RGBColor(0x3c, 0x3a, 0x38)
GREY  = RGBColor(0x6f, 0x6d, 0x69)
FAINT = RGBColor(0xa8, 0xa6, 0xa2)
LINEC = RGBColor(0xdc, 0xdb, 0xd7)
HEAD = "GT America"
BODY = "GT America Light"

prs = Presentation()
SW, SH = Inches(13.333), Inches(7.5)
prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]
ML = Inches(0.7); CW = SW - 2 * ML
_pageno = [0]


def slide():
    return prs.slides.add_slide(BLANK)


def tb(s, left, top, width, height, paras, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = s.shapes.add_textbox(left, top, width, height); tf = box.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "line" in para: p.line_spacing = para["line"]
        if "after" in para: p.space_after = Pt(para["after"])
        for run in para["runs"]:
            r = p.add_run(); r.text = run["t"]; f = r.font
            f.name = run.get("font", BODY); f.size = Pt(run["size"])
            f.color.rgb = run.get("color", INK)
            f.bold = run.get("bold", False); f.italic = run.get("italic", False)
            if run.get("track"):
                r._r.get_or_add_rPr().set("spc", str(int(run["track"] * 100)))
    return box


def rule(s, top, left=ML, width=CW):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.3))
    sp.fill.solid(); sp.fill.fore_color.rgb = INK; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def header(s, label, route):
    tb(s, ML, Inches(0.46), Inches(9), Inches(0.32),
       [{"runs": [{"t": label.upper(), "font": HEAD, "size": 11.5, "color": INK, "bold": True, "track": 1.7}]}])
    tb(s, Inches(4.6), Inches(0.49), CW - Inches(3.9), Inches(0.3),
       [{"runs": [{"t": route, "font": BODY, "size": 10, "color": FAINT, "italic": True, "track": 0.6}], "align": PP_ALIGN.RIGHT}])
    rule(s, Inches(0.84))


def footer(s, *, wm=True):
    _pageno[0] += 1
    if wm:
        tb(s, ML, Inches(7.12), Inches(3), Inches(0.25),
           [{"runs": [{"t": "PERENNIAL", "font": HEAD, "size": 8.5, "color": FAINT, "bold": True, "track": 1.6}]}])
    tb(s, SW - ML - Inches(1.2), Inches(7.12), Inches(1.2), Inches(0.25),
       [{"runs": [{"t": f"{_pageno[0]:02d}", "font": BODY, "size": 9, "color": FAINT, "track": 1.0}], "align": PP_ALIGN.RIGHT}])


def cpic(s, img, top, height, *, border=True):
    pic = s.shapes.add_picture(os.path.join(SHOTS, img), 0, top, height=height)
    pic.left = int((SW - pic.width) / 2)
    if border:
        pic.line.color.rgb = LINEC; pic.line.width = Pt(0.75)
    return pic


def caption(s, text):
    tb(s, Inches(1.6), Inches(6.12), Inches(10.13), Inches(0.95),
       [{"runs": [{"t": text, "font": BODY, "size": 13, "color": SOFT}], "align": PP_ALIGN.CENTER, "line": 1.32}],
       anchor=MSO_ANCHOR.TOP)


def content(label, route, img, cap):
    s = slide(); header(s, label, route)
    cpic(s, img, Inches(1.0), Inches(5.0)); caption(s, cap); footer(s)


# ---- 1 · TITLE (split) ----
s = slide()
pic = s.shapes.add_picture(os.path.join(SHOTS, "home.png"), Inches(0.55), 0, width=Inches(6.55))
pic.top = int((SH - pic.height) / 2)
pic.line.color.rgb = LINEC; pic.line.width = Pt(0.75)
tx = Inches(7.5); tw = Inches(5.25)
tb(s, tx, Inches(1.5), tw, Inches(0.6),
   [{"runs": [{"t": "Perennial", "font": HEAD, "size": 33, "color": INK, "bold": True, "track": -0.3}]}])
tb(s, tx, Inches(2.22), tw, Inches(0.7),
   [{"runs": [{"t": "Studio-management software for designers and makers — designed and built end to end.",
               "font": BODY, "size": 14, "color": GREY, "italic": True}], "line": 1.25}])
tb(s, tx, Inches(3.2), tw, Inches(2.85),
   [{"runs": [
       {"t": "Perennial helps artists and designers turn their craft into a business, so they can spend more time making the work because it finally pays for itself. ",
        "font": BODY, "size": 13, "color": SOFT},
       {"t": "It gives an independent practice the few tangible pieces it actually needs — clients, projects, money, and visibility — without assuming you speak the language of operations.",
        "font": HEAD, "size": 13, "color": INK, "bold": True}], "line": 1.42}])
tb(s, tx, Inches(6.2), tw, Inches(0.3),
   [{"runs": [{"t": "PRODUCT WALKTHROUGH   ·   APP.PERENNIAL.DESIGN   ·   ELLIOTT ROSENBERG",
               "font": BODY, "size": 9.5, "color": FAINT, "track": 1.0}]}])
footer(s)

# ---- 2..14 content ----
CONTENT = [
    ("Home — Dashboard", "/home", "home.png",
     "The home screen pulls together data from every other module into a single view. Tasks coming due, invoices outstanding, active projects, and recent notes all surface here, so the first thing you see each day is an accurate read on the studio — without opening anything else."),
    ("Projects — Board", "/projects", "projects.png",
     "Projects are arranged on a board by status, and each card shows its type, client, value, and progress. The progress bar is calculated from the project’s underlying tasks, so the board stays current on its own as work gets checked off."),
    ("Projects — A project’s tasks", "/projects · detail", "project-scrim-tasks.png",
     "Opening a project reveals a panel that holds everything about it in one place. The Tasks tab is a checklist with due dates and priorities, and the time logged against these tasks feeds directly into the studio’s billable hours and, eventually, an invoice."),
    ("Projects — A project’s notes", "/projects · detail", "project-scrim-notes.png",
     "The same panel keeps the project’s notes beside its tasks, files, and people. Research, creative direction, client constraints, and scope all live with the work they belong to, rather than scattered across a separate notes app."),
    ("Network — Contacts", "/network", "network-contacts.png",
     "The Network module is a CRM for everyone the studio works with — clients, collaborators, press, and vendors. Each person is linked to their organization and to the projects they’re part of, so a contact, a company, and a piece of work are never more than a click apart."),
    ("Network — Leads", "/network · leads", "network-leads.png",
     "Potential work is tracked as leads that move through stages, from first contact to qualified. New business, press, and stockists all run through the same pipeline, and one person can be a client, a lead, and a press contact at once without duplicating their record."),
    ("Network — Organizations", "/network · organizations", "network-orgs.png",
     "Companies are tracked separately from people and tagged by type — brands, publications, galleries, fairs, and vendors. Opening an organization shows the people who work there and the history the studio has built with them."),
    ("Outreach — Pipelines", "/outreach", "outreach.png",
     "Outreach turns those relationships into visual pipelines for pursuing press, galleries, stockists, and new business. Each target moves through stages as it progresses, and the outcomes recorded here connect back to the contacts and coverage tracked elsewhere in the app."),
    ("Finance — Invoices", "/finance · invoices", "invoices-draft.png",
     "Invoices are built from data the app already holds — the client, the project, and the rate carry over automatically. The selected invoice opens its full editor, where line items can be drawn straight from logged time before the invoice is sent."),
    ("Finance — A sent invoice", "/finance · invoices", "invoice-sent.png",
     "Once sent, an invoice becomes a clean, branded document carrying the studio’s name and colors. Clients pay online through Stripe, with funds landing directly in the studio’s own connected account — and the payment is later matched against the bank feed."),
    ("Finance — Banking", "/finance · banking", "banking.png",
     "A live bank feed connected through Plaid brings every transaction in and categorizes it automatically. Incoming payments are matched back to the invoice that produced them, so reconciling what was billed against what actually landed takes a single click."),
    ("Presence — Opportunities", "/presence · opportunities", "presence-opps.png",
     "The Opportunities feed is a curated, continuously updated list of fairs, open calls, grants, and awards relevant to independent studios. Items can be filtered by type and saved, where they connect into the calendar and outreach for follow-up."),
    ("Ash — Assistant", "global", "ash.png",
     "Ash is built into every screen and already understands the studio — your clients, your projects, how your money moves, even how you write. Rather than sending you elsewhere with generic answers, it can do the task itself: draft the client update, build a cost list, or surface what to prioritize, using the context the rest of the app already holds."),
]
for label, route, img, cap in CONTENT:
    content(label, route, img, cap)

# ---- 15 · CLOSE ----
s = slide()
tb(s, Inches(1.5), Inches(2.45), Inches(10.33), Inches(0.9),
   [{"runs": [{"t": "Perennial", "font": HEAD, "size": 44, "color": INK, "bold": True, "track": -0.4}], "align": PP_ALIGN.CENTER}])
tb(s, Inches(1.5), Inches(3.5), Inches(10.33), Inches(0.4),
   [{"runs": [{"t": "Studio-management software for designers and makers.", "font": BODY, "size": 16, "color": GREY, "italic": True}], "align": PP_ALIGN.CENTER}])
rule(s, Inches(4.25), left=Inches(5.47), width=Inches(2.4))
tb(s, Inches(1.5), Inches(4.55), Inches(10.33), Inches(0.35),
   [{"runs": [{"t": "app.perennial.design", "font": HEAD, "size": 14, "color": INK, "bold": True}], "align": PP_ALIGN.CENTER}])
tb(s, Inches(1.5), Inches(5.0), Inches(10.33), Inches(0.3),
   [{"runs": [{"t": "DESIGNED & BUILT END TO END BY ELLIOTT ROSENBERG", "font": BODY, "size": 10, "color": FAINT, "track": 1.2}], "align": PP_ALIGN.CENTER}])
footer(s, wm=False)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
